"""Office document parsers: .docx (python-docx), .pptx (python-pptx), .xlsx (MarkItDown)."""

from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation

from ._models import (
    ParsedDocument,
    _INTRO_HEADING,
    _MARKITDOWN,
    _sections_to_document,
    _split_by_heading,
)


def _parse_docx(path: Path) -> ParsedDocument:
    """Parse Word .docx using python-docx; split on Heading 2 style."""
    doc = Document(path)
    title = path.stem
    source_url = ""
    raw_sections: list[tuple[str, str]] = []
    current_heading = _INTRO_HEADING
    current_paras: list[str] = []

    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        if style == "Title" and title == path.stem:
            title = text
        elif style == "Heading 2":
            if current_paras:
                raw_sections.append((current_heading, "\n\n".join(current_paras)))
            current_heading = text
            current_paras = []
        else:
            current_paras.append(text)

    if current_paras:
        raw_sections.append((current_heading, "\n\n".join(current_paras)))
    return _sections_to_document(title, source_url, raw_sections)


def _parse_pptx(path: Path) -> ParsedDocument:
    """Parse PowerPoint .pptx; each slide becomes one L2 section."""
    prs = Presentation(path)
    title = path.stem
    source_url = ""
    raw_sections: list[tuple[str, str]] = []

    for i, slide in enumerate(prs.slides):
        slide_heading = f"Slide {i + 1}"
        texts: list[str] = []

        for shape_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if not text:
                    continue
                if i == 0 and shape_idx == 0 and para_idx == 0 and title == path.stem:
                    title = text
                if shape_idx == 0 and para_idx == 0:
                    slide_heading = text
                else:
                    texts.append(text)

        if texts:
            raw_sections.append((slide_heading, "\n\n".join(texts)))

    return _sections_to_document(title, source_url, raw_sections)


def _parse_xlsx(path: Path) -> ParsedDocument:
    """Parse Excel .xlsx using MarkItDown; each sheet becomes one L2 section."""
    result = _MARKITDOWN.convert(str(path))
    markdown = result.text_content or ""
    raw_sections = _split_by_heading(markdown, pattern=r"^##\s+(.+)$")
    if not raw_sections or all(not t.strip() for _, t in raw_sections):
        raw_sections = [(_INTRO_HEADING, markdown.strip())]
    return _sections_to_document(path.stem, "", raw_sections)
